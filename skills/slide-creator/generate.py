#!/usr/bin/env python3
"""
Slide Creator - EC戦略資料生成スクリプト
テンプレートPPTXをベースに新しい資料を生成する
"""

import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

SKILL_DIR = Path(__file__).parent
DEFAULT_TEMPLATE = SKILL_DIR / "templates" / "ec-strategy.pptx"


def load_template(template_path: Path) -> Presentation:
    """テンプレートPPTXを読み込む"""
    return Presentation(str(template_path))


def replace_text_in_shape(shape, replacements: dict):
    """シェイプ内のテキストを置換する"""
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)


def replace_text_in_slide(slide, replacements: dict):
    """スライド内の全テキストを置換する"""
    for shape in slide.shapes:
        replace_text_in_shape(shape, replacements)
        
        # テーブル内のテキストも置換
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old, new in replacements.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)


def create_presentation_from_config(config: dict, template_path: Path, output_path: Path):
    """設定ファイルから資料を生成する"""
    prs = load_template(template_path)
    
    # グローバル置換（全スライドに適用）
    global_replacements = {
        "アントレスクエア": config.get("project_name", "プロジェクト名"),
        "株式会社アントレックス": config.get("company", "株式会社〇〇"),
        "WCA加藤": config.get("author", "担当者"),
    }
    
    for slide in prs.slides:
        replace_text_in_slide(slide, global_replacements)
    
    # スライドごとの置換（configにslidesがあれば）
    if "slides" in config:
        for i, slide_config in enumerate(config["slides"]):
            if i < len(prs.slides):
                slide = prs.slides[i]
                if "replacements" in slide_config:
                    replace_text_in_slide(slide, slide_config["replacements"])
    
    # 保存
    prs.save(str(output_path))
    print(f"✅ 生成完了: {output_path}")
    return output_path


def create_simple_presentation(
    project_name: str,
    company: str,
    author: str,
    template_path: Path,
    output_path: Path
):
    """シンプルな置換のみで資料を生成する"""
    config = {
        "project_name": project_name,
        "company": company,
        "author": author,
    }
    return create_presentation_from_config(config, template_path, output_path)


def main():
    parser = argparse.ArgumentParser(description="EC戦略資料を生成する")
    parser.add_argument("--template", "-t", type=Path, default=DEFAULT_TEMPLATE,
                        help="テンプレートPPTXファイル")
    parser.add_argument("--output", "-o", type=Path, required=True,
                        help="出力ファイルパス")
    parser.add_argument("--config", "-c", type=Path,
                        help="設定JSONファイル")
    parser.add_argument("--project", "-p", type=str,
                        help="プロジェクト名")
    parser.add_argument("--company", type=str,
                        help="会社名")
    parser.add_argument("--author", "-a", type=str,
                        help="担当者名")
    
    args = parser.parse_args()
    
    if args.config:
        # 設定ファイルから生成
        with open(args.config, "r", encoding="utf-8") as f:
            config = json.load(f)
        create_presentation_from_config(config, args.template, args.output)
    elif args.project:
        # コマンドライン引数から生成
        create_simple_presentation(
            project_name=args.project,
            company=args.company or "株式会社〇〇",
            author=args.author or "担当者",
            template_path=args.template,
            output_path=args.output
        )
    else:
        print("❌ --config または --project を指定してください")
        sys.exit(1)


if __name__ == "__main__":
    main()
